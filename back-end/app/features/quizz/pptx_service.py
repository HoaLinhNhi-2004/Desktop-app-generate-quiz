"""PPTX Service - Extract text from PowerPoint decks.

Slide decks are the single most common study material students hand over, and
`python-pptx` was already a dependency (Smart Import used it for a 500-character
classification snippet). What was missing was an extractor: Smart Import accepted
`.pptx`, created the record, and then `document_processor._extract_file` raised
`Unsupported file type: pptx` because it had no branch for it.

Slide boundaries are kept as `--- SLIDE N ---` markers, mirroring the
`--- TRANG N ---` convention `pdf_service` uses, so downstream chunking and the
source-page heatmap see the same shape for both.
"""

import os
import logging

logger = logging.getLogger(__name__)


def _shape_text(shape) -> list[str]:
    """Text from one shape, including tables and grouped shapes."""
    out: list[str] = []

    if shape.shape_type is not None and getattr(shape, "shapes", None) is not None:
        # Grouped shapes nest arbitrarily deep.
        for child in shape.shapes:
            out.extend(_shape_text(child))
        return out

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                out.append(" | ".join(cells))
        return out

    if getattr(shape, "has_text_frame", False):
        for para in shape.text_frame.paragraphs:
            line = "".join(run.text for run in para.runs).strip()
            if line:
                out.append(line)

    return out


def extract_text_from_pptx(pptx_path: str) -> str:
    """Extract text from a .pptx file, one annotated block per slide.

    Speaker notes are included: for lecture decks the slide itself is often only
    a few bullet points and the actual explanation lives in the notes.
    """
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    from pptx import Presentation

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        logger.error("Error reading PPTX %s: %s", os.path.basename(pptx_path), e)
        raise ValueError(
            "Không đọc được file PowerPoint. File có thể hỏng hoặc là định dạng "
            ".ppt cũ — hãy lưu lại thành .pptx."
        ) from e

    blocks: list[str] = []
    for index, slide in enumerate(prs.slides, 1):
        lines: list[str] = []
        for shape in slide.shapes:
            lines.extend(_shape_text(shape))

        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                lines.append(f"[Ghi chú] {notes}")

        if lines:
            blocks.append(f"--- SLIDE {index} ---\n" + "\n".join(lines))

    combined = "\n\n".join(blocks)
    logger.info(
        "Extracted %d chars from %d slide(s) in %s",
        len(combined), len(blocks), os.path.basename(pptx_path),
    )
    return combined
