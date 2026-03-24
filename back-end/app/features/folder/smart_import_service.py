"""
Smart Folder Import Service — Scan a local directory, extract text snippets,
use Gemini AI to categorize files into learning-material folders, and import them.

Pipeline:
  1. Walk the directory recursively, filter supported file types
  2. Smart-extract representative text from each file (head + tail + densest page)
  3. Batch files (15 per request) → send to Gemini for categorization
  4. Match AI-suggested folders to existing DB folders or create new ones
  5. TRASH files → "Review Later" folder (safe, not deleted)
  6. Copy files to upload storage and create UploadedFileRecord entries
  7. Progress tracking via HTTP polling with pause/cancel support
"""

import os
import re
import uuid
import json
import time
import shutil
import hashlib
import logging
import threading
from datetime import datetime, timezone

from app.db import db
from app.features.folder.models import Folder
from app.features.upload.models import UploadedFileRecord

logger = logging.getLogger(__name__)

# Supported file extensions for import
SUPPORTED_EXTENSIONS = {"pdf", "docx", "doc", "txt", "pptx"}

# Files to always skip (system / hidden files)
SKIP_FILES = {
    "desktop.ini", "thumbs.db", ".ds_store", ".gitignore",
    ".gitkeep", "icon.ico", "folder.ico",
}

BATCH_SIZE = 15  # files per Gemini API call
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50 MB — skip files larger than this
REVIEW_FOLDER_NAME = "Review Later"  # folder for AI-rejected files
MAX_RETRY_WAIT = 60  # seconds to wait before retry on rate limit

# Folder colors for auto-created folders (cycle through these)
AUTO_FOLDER_COLORS = [
    "hsl(262 83% 58%)",
    "hsl(217 91% 60%)",
    "hsl(142 71% 45%)",
    "hsl(38 92% 50%)",
    "hsl(346 87% 58%)",
    "hsl(187 88% 40%)",
]

# ── Global in-memory progress store (keyed by job_id) ────────────────────────
_import_jobs: dict[str, dict] = {}
_jobs_lock = threading.Lock()


def get_job(job_id: str) -> dict | None:
    with _jobs_lock:
        return _import_jobs.get(job_id)


def _update_job(job_id: str, **kwargs):
    with _jobs_lock:
        if job_id in _import_jobs:
            _import_jobs[job_id].update(kwargs)


def _is_job_cancelled(job_id: str) -> bool:
    with _jobs_lock:
        job = _import_jobs.get(job_id)
        return job is not None and job.get("cancelRequested", False)


def _is_job_paused(job_id: str) -> bool:
    with _jobs_lock:
        job = _import_jobs.get(job_id)
        return job is not None and job.get("paused", False)


def cancel_job(job_id: str) -> bool:
    """Request cancellation of a running job."""
    with _jobs_lock:
        if job_id in _import_jobs:
            _import_jobs[job_id]["cancelRequested"] = True
            _import_jobs[job_id]["paused"] = False  # unpause so loop can exit
            return True
    return False


def pause_job(job_id: str) -> bool:
    """Pause a running job."""
    with _jobs_lock:
        if job_id in _import_jobs and _import_jobs[job_id]["status"] not in ("completed", "error", "cancelled"):
            _import_jobs[job_id]["paused"] = True
            return True
    return False


def resume_job(job_id: str) -> bool:
    """Resume a paused job."""
    with _jobs_lock:
        if job_id in _import_jobs:
            _import_jobs[job_id]["paused"] = False
            return True
    return False


def _add_file_status(job_id: str, filename: str, status: str, folder_name: str = "", reason: str = ""):
    with _jobs_lock:
        if job_id not in _import_jobs:
            return
        _import_jobs[job_id]["files"].append({
            "name": filename,
            "status": status,        # scanning | categorizing | done | skipped | error | review
            "folderName": folder_name,
            "reason": reason,
        })
        # Update counters
        if status == "done":
            _import_jobs[job_id]["completed"] += 1
        elif status in ("skipped", "error"):
            _import_jobs[job_id]["skipped"] += 1
        elif status == "review":
            _import_jobs[job_id]["reviewCount"] = _import_jobs[job_id].get("reviewCount", 0) + 1


def _update_file_status(job_id: str, filename: str, status: str, folder_name: str = "", reason: str = ""):
    with _jobs_lock:
        if job_id not in _import_jobs:
            return
        for f in _import_jobs[job_id]["files"]:
            if f["name"] == filename:
                f["status"] = status
                if folder_name:
                    f["folderName"] = folder_name
                if reason:
                    f["reason"] = reason
                break
        if status == "done":
            _import_jobs[job_id]["completed"] += 1
        elif status in ("skipped", "error"):
            _import_jobs[job_id]["skipped"] += 1
        elif status == "review":
            _import_jobs[job_id]["reviewCount"] = _import_jobs[job_id].get("reviewCount", 0) + 1


# ── Path Validation ──────────────────────────────────────────────────────────

def _validate_dir_path(dir_path: str) -> str:
    """Validate and normalize a directory path to prevent directory traversal."""
    normalized = os.path.normpath(os.path.abspath(dir_path))
    if not os.path.isdir(normalized):
        raise ValueError(f"Directory not found: {dir_path}")
    # Block paths that contain suspicious traversal patterns
    if ".." in os.path.relpath(normalized, os.path.splitdrive(normalized)[0] + os.sep):
        raise ValueError("Invalid directory path")
    return normalized


# ── Step 1: Scan directory ───────────────────────────────────────────────────

def scan_directory(dir_path: str) -> list[dict]:
    """
    Walk the directory tree and return a list of supported files.
    Skips files > 50MB and system files.
    Each item: {"path": str, "name": str, "ext": str, "rel_dir": str, "size": int}
    """
    result = []
    dir_path = os.path.normpath(dir_path)

    for root, _dirs, files in os.walk(dir_path):
        for fname in files:
            if fname.lower() in SKIP_FILES:
                continue
            ext = fname.rsplit(".", 1)[-1].lower() if "." in fname else ""
            if ext not in SUPPORTED_EXTENSIONS:
                continue
            full_path = os.path.join(root, fname)
            try:
                fsize = os.path.getsize(full_path)
            except OSError:
                continue
            if fsize == 0:
                continue
            if fsize > MAX_FILE_SIZE:
                logger.info("Skipping large file (%d MB): %s", fsize // (1024 * 1024), fname)
                continue
            # Relative subdirectory from the import root
            rel_dir = os.path.relpath(root, dir_path)
            if rel_dir == ".":
                rel_dir = ""
            result.append({
                "path": full_path,
                "name": fname,
                "ext": ext,
                "rel_dir": rel_dir,
                "size": fsize,
            })

    return result


# ── Deduplication ────────────────────────────────────────────────────────────

def _compute_file_hash(file_path: str) -> str:
    """Compute a fast partial hash (first + last 4KB) for deduplication."""
    h = hashlib.md5()
    try:
        with open(file_path, "rb") as f:
            head = f.read(4096)
            h.update(head)
            f.seek(0, 2)  # end of file
            fsize = f.tell()
            if fsize > 4096:
                f.seek(max(0, fsize - 4096))
                h.update(f.read(4096))
            h.update(str(fsize).encode())
    except OSError:
        return ""
    return h.hexdigest()


def _get_existing_file_hashes() -> set[str]:
    """Get hashes of files already in uploads folder."""
    hashes = set()
    records = UploadedFileRecord.query.with_entities(UploadedFileRecord.stored_path).all()
    for (stored_path,) in records:
        if stored_path and os.path.isfile(stored_path):
            h = _compute_file_hash(stored_path)
            if h:
                hashes.add(h)
    return hashes


# ── Step 2: Smart text extraction ────────────────────────────────────────────

def _extract_snippet(file_path: str, ext: str, max_chars: int = 500) -> str:
    """
    Smart extraction: grab text from head, tail, and the densest page
    to give AI maximum context for classification.
    """
    try:
        if ext == "pdf":
            return _extract_pdf_smart(file_path, max_chars)
        elif ext in ("docx", "doc"):
            return _extract_docx_smart(file_path, max_chars)
        elif ext == "txt":
            return _extract_txt_smart(file_path, max_chars)
        elif ext == "pptx":
            return _extract_pptx_smart(file_path, max_chars)
    except Exception as e:
        logger.warning("Snippet extraction failed for %s: %s", file_path, e)
    return ""


def _extract_pdf_smart(path: str, max_chars: int) -> str:
    """Extract from first page, last page, and densest page of a PDF."""
    import fitz  # PyMuPDF
    try:
        doc = fitz.open(path)
    except Exception:
        return ""

    if doc.page_count == 0:
        doc.close()
        return ""

    per_section = max_chars // 3  # ~166 chars each for head/tail/dense

    # Head: first page text
    head_text = doc[0].get_text()[:per_section]

    # Tail: last page text
    tail_text = ""
    if doc.page_count > 1:
        tail_text = doc[-1].get_text()[:per_section]

    # Densest page: the page with the most text (skip first & last)
    dense_text = ""
    if doc.page_count > 2:
        max_len = 0
        densest_page = None
        for i in range(1, doc.page_count - 1):
            page_text = doc[i].get_text()
            if len(page_text) > max_len:
                max_len = len(page_text)
                densest_page = page_text
        if densest_page:
            dense_text = densest_page[:per_section]

    doc.close()

    parts = [p.strip() for p in [head_text, dense_text, tail_text] if p.strip()]
    return " [...] ".join(parts)[:max_chars].strip()


def _extract_docx_smart(path: str, max_chars: int) -> str:
    """Extract head + tail from DOCX."""
    from app.features.quizz.docx_service import extract_text_from_docx
    full_text = extract_text_from_docx(path)
    if not full_text:
        return ""
    per_section = max_chars // 2
    head = full_text[:per_section]
    tail = full_text[-per_section:] if len(full_text) > per_section else ""
    if tail and tail != head:
        return (head.strip() + " [...] " + tail.strip())[:max_chars]
    return head[:max_chars].strip()


def _extract_txt_smart(path: str, max_chars: int) -> str:
    """Extract head + tail from TXT."""
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()
    if not content:
        return ""
    per_section = max_chars // 2
    head = content[:per_section]
    tail = content[-per_section:] if len(content) > per_section else ""
    if tail and tail != head:
        return (head.strip() + " [...] " + tail.strip())[:max_chars]
    return head[:max_chars].strip()


def _extract_pptx_smart(path: str, max_chars: int) -> str:
    """Extract from first slides and last slides of PPTX."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = list(prs.slides)
        if not slides:
            return ""

        per_section = max_chars // 2

        # Head: first few slides
        head_text = ""
        for slide in slides[:3]:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    head_text += shape.text + " "
                if len(head_text) >= per_section:
                    break
            if len(head_text) >= per_section:
                break

        # Tail: last few slides
        tail_text = ""
        for slide in reversed(slides[-3:]):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    tail_text += shape.text + " "
                if len(tail_text) >= per_section:
                    break
            if len(tail_text) >= per_section:
                break

        parts = [p.strip() for p in [head_text, tail_text] if p.strip()]
        return " [...] ".join(parts)[:max_chars].strip()
    except ImportError:
        return ""


# ── Step 3: Gemini AI Categorization (Batched, with improved prompt) ─────────

def _normalize_folder_name(name: str) -> str:
    """Standardize folder names: title case, strip excess whitespace, limit length."""
    name = name.strip()
    if not name:
        return "Uncategorized"
    # Title case
    name = name.title()
    # Remove excess internal whitespace
    name = re.sub(r"\s+", " ", name)
    # Limit to 50 chars
    if len(name) > 50:
        name = name[:50].rsplit(" ", 1)[0]
    return name


def _categorize_batch(
    file_infos: list[dict],
    existing_folders: list[str],
    api_key: str,
    model_chain: list[str],
    job_id: str = "",
) -> dict[str, str]:
    """
    Send a batch of files to Gemini and return a mapping: filename → folder name.
    Uses 'TRASH' for non-educational files. Includes few-shot examples and
    structured output instructions for better accuracy.
    """
    import google.generativeai as genai

    # Build the file list for the prompt
    file_list_text = ""
    for i, fi in enumerate(file_infos, 1):
        snippet = fi.get("snippet", "")
        file_list_text += (
            f"\n{i}. File: \"{fi['name']}\"\n"
            f"   Subdirectory: \"{fi.get('rel_dir', '')}\"\n"
            f"   Content excerpt: \"{snippet[:400]}\"\n"
        )

    existing_str = ", ".join(f'"{f}"' for f in existing_folders) if existing_folders else "(none)"

    prompt = (
        "You are an expert document classification AI specializing in educational materials.\n"
        "Below is a list of files with their names, subdirectories, and content excerpts.\n\n"
        "TASK:\n"
        "1. Determine if each file is educational/study material (textbooks, lecture notes, exercises, "
        "research papers, course material, lab reports, etc.).\n"
        "2. For educational files, assign them to the MOST APPROPRIATE folder from the existing folders "
        f"list: [{existing_str}]. If no existing folder matches well, suggest a SHORT new folder name "
        "(2-4 words, Title Case, e.g. 'Mathematics', 'Computer Science', 'Vietnamese History').\n"
        "3. For NON-educational files (receipts, invoices, personal documents, fiction, code files, etc.), "
        "assign folder = 'TRASH'.\n\n"
        "RULES FOR FOLDER NAMES:\n"
        "- Use Title Case (capitalize first letter of each word)\n"
        "- Keep names 2-4 words maximum\n"
        "- Be consistent: don't create 'Math' and 'Mathematics' separately\n"
        "- Prefer existing folder names when content matches\n\n"
        "FEW-SHOT EXAMPLES:\n"
        '- File "GiaiTich1_Chuong3.pdf", content about derivatives and integrals → folder: "Mathematics"\n'
        '- File "Hoa_don_dien_T3.pdf", content about electricity bill payment → folder: "TRASH"\n'
        '- File "Slide_OOP_Java.pptx", content about classes and inheritance → folder: "Computer Science"\n'
        '- File "TieuLuan_LichSu.docx", content about Vietnamese history essay → folder: "Vietnamese History"\n'
        '- File "Novel_ConanChap1.pdf", content about detective fiction story → folder: "TRASH"\n\n'
        "RESPOND with ONLY a valid JSON array (no markdown fences, no explanation):\n"
        '[{"file": "exact_filename.pdf", "folder": "Folder Name"}, ...]\n\n'
        f"FILES:{file_list_text}"
    )

    genai.configure(api_key=api_key)

    last_err = None
    for model_name in model_chain:
        try:
            model = genai.GenerativeModel(model_name)
            response = model.generate_content(
                prompt,
                generation_config=genai.types.GenerationConfig(
                    temperature=0.1,  # lower for more consistent classification
                    max_output_tokens=4096,
                ),
            )
            raw = (response.text or "").strip()

            # Parse JSON response
            cleaned = re.sub(r"```(?:json)?\s*", "", raw).strip()
            cleaned = re.sub(r"```\s*$", "", cleaned).strip()

            parsed = json.loads(cleaned)
            if isinstance(parsed, list):
                mapping = {}
                for item in parsed:
                    fname = item.get("file", "")
                    folder = item.get("folder", "TRASH")
                    # Normalize folder names (except TRASH)
                    if folder.upper() != "TRASH":
                        folder = _normalize_folder_name(folder)
                    mapping[fname] = folder
                return mapping
        except Exception as e:
            err_str = str(e).lower()
            is_429 = "429" in err_str or "resource_exhausted" in err_str
            if is_429:
                logger.warning("Model %s rate limited, trying next model...", model_name)
                # Update job status so frontend can see rate limit message
                if job_id:
                    _update_job(job_id, rateLimitInfo=f"Rate limited on {model_name}, trying fallback...")
                last_err = e
                continue
            logger.error("Gemini categorization error with %s: %s", model_name, e)
            last_err = e
            continue

    # All models exhausted — try retry with backoff
    if last_err and job_id:
        err_str = str(last_err).lower()
        is_429 = "429" in err_str or "resource_exhausted" in err_str
        if is_429:
            _update_job(job_id, rateLimitInfo="All models rate limited. Waiting 60s before retry...")
            logger.info("All models rate limited. Waiting %ds before retry...", MAX_RETRY_WAIT)
            # Wait with cancel check
            for _ in range(MAX_RETRY_WAIT):
                if _is_job_cancelled(job_id):
                    return {}
                time.sleep(1)
            _update_job(job_id, rateLimitInfo="")
            # Retry once with the first model
            try:
                import google.generativeai as genai
                genai.configure(api_key=api_key)
                model = genai.GenerativeModel(model_chain[0])
                response = model.generate_content(
                    prompt,
                    generation_config=genai.types.GenerationConfig(
                        temperature=0.1,
                        max_output_tokens=4096,
                    ),
                )
                raw = (response.text or "").strip()
                cleaned = re.sub(r"```(?:json)?\s*", "", raw).strip()
                cleaned = re.sub(r"```\s*$", "", cleaned).strip()
                parsed = json.loads(cleaned)
                if isinstance(parsed, list):
                    mapping = {}
                    for item in parsed:
                        fname = item.get("file", "")
                        folder = item.get("folder", "TRASH")
                        if folder.upper() != "TRASH":
                            folder = _normalize_folder_name(folder)
                        mapping[fname] = folder
                    return mapping
            except Exception as retry_err:
                logger.error("Retry also failed: %s", retry_err)

    logger.error("All models failed for categorization. Last error: %s", last_err)
    # Fallback: use subdirectory structure
    return {fi["name"]: _normalize_folder_name(fi.get("rel_dir", "") or "Uncategorized") for fi in file_infos}


# ── Step 4: Main import pipeline (runs in background thread) ────────────────

def start_import_job(dir_path: str, app) -> str:
    """
    Start a background import job. Returns a job_id for progress tracking.
    """
    # Validate path
    dir_path = _validate_dir_path(dir_path)

    job_id = str(uuid.uuid4())

    with _jobs_lock:
        _import_jobs[job_id] = {
            "id": job_id,
            "status": "scanning",  # scanning | categorizing | importing | completed | error | cancelled
            "dirPath": dir_path,
            "totalFiles": 0,
            "completed": 0,
            "skipped": 0,
            "reviewCount": 0,
            "files": [],
            "createdFolders": [],
            "error": None,
            "startedAt": datetime.now(timezone.utc).isoformat(),
            "paused": False,
            "cancelRequested": False,
            "rateLimitInfo": "",
        }

    t = threading.Thread(
        target=_run_import_pipeline,
        args=(job_id, dir_path, app),
        daemon=True,
    )
    t.start()

    return job_id


def _wait_while_paused(job_id: str):
    """Block the pipeline while the job is paused. Returns True if cancelled."""
    while _is_job_paused(job_id):
        if _is_job_cancelled(job_id):
            return True
        time.sleep(0.5)
    return _is_job_cancelled(job_id)


def _run_import_pipeline(job_id: str, dir_path: str, app):
    """The main import pipeline, running in a background thread."""
    with app.app_context():
        try:
            # Step 1: Scan
            _update_job(job_id, status="scanning")
            all_files = scan_directory(dir_path)

            if not all_files:
                _update_job(job_id, status="completed", totalFiles=0)
                return

            _update_job(job_id, totalFiles=len(all_files))

            # Register all files as "pending"
            for fi in all_files:
                _add_file_status(job_id, fi["name"], "scanning")

            # Check cancel
            if _is_job_cancelled(job_id):
                _update_job(job_id, status="cancelled")
                return

            # Step 2: Extract snippets (Smart Deep Scan)
            _update_job(job_id, status="categorizing")

            # Deduplication: get existing file hashes
            existing_hashes = _get_existing_file_hashes()

            files_to_categorize = []
            for fi in all_files:
                if _wait_while_paused(job_id):
                    _update_job(job_id, status="cancelled")
                    return

                # Dedup check
                file_hash = _compute_file_hash(fi["path"])
                if file_hash and file_hash in existing_hashes:
                    _update_file_status(job_id, fi["name"], "skipped", reason="Already imported (duplicate)")
                    continue

                fi["hash"] = file_hash
                snippet = _extract_snippet(fi["path"], fi["ext"])
                fi["snippet"] = snippet

                if not snippet.strip():
                    _update_file_status(job_id, fi["name"], "skipped", reason="Could not extract text (may be scanned PDF or password-protected)")
                    continue

                files_to_categorize.append(fi)

            if not files_to_categorize:
                _update_job(job_id, status="completed")
                return

            # Get API key from the key manager
            from app.features.api_keys.key_manager import get_optimal_key
            key_obj = get_optimal_key()
            if not key_obj:
                _update_job(job_id, status="error", error="No API key available. Please add a Gemini API key in Settings.")
                return

            api_key = key_obj.key
            from flask import current_app
            model_chain = current_app.config.get("GEMINI_FALLBACK_CHAIN", ["gemini-2.5-flash", "gemini-2.5-flash-lite", "gemini-2.0-flash"])

            # Get existing folder names
            existing_folders = [f.name for f in Folder.query.all()]

            # Step 3: Batch categorization
            full_mapping: dict[str, str] = {}
            batches = [files_to_categorize[i:i + BATCH_SIZE] for i in range(0, len(files_to_categorize), BATCH_SIZE)]

            for batch_idx, batch in enumerate(batches):
                if _wait_while_paused(job_id):
                    _update_job(job_id, status="cancelled")
                    return

                logger.info("Categorizing batch %d/%d (%d files)", batch_idx + 1, len(batches), len(batch))

                # Update status for files in this batch
                for fi in batch:
                    _update_file_status(job_id, fi["name"], "categorizing")

                batch_mapping = _categorize_batch(batch, existing_folders, api_key, model_chain, job_id)

                if _is_job_cancelled(job_id):
                    _update_job(job_id, status="cancelled")
                    return

                full_mapping.update(batch_mapping)

                # Update existing_folders with any new names from this batch
                for folder_name in batch_mapping.values():
                    if folder_name.upper() != "TRASH" and folder_name not in existing_folders:
                        existing_folders.append(folder_name)

                # Rate-limit safety delay between batches
                if batch_idx < len(batches) - 1:
                    time.sleep(4)

            # Step 4: Import files — create folders + copy files + DB records
            _update_job(job_id, status="importing", rateLimitInfo="")

            upload_folder = current_app.config.get("UPLOAD_FOLDER", "uploads")
            os.makedirs(upload_folder, exist_ok=True)

            # Build folder cache (name → Folder object)
            folder_cache: dict[str, Folder] = {}
            for f in Folder.query.all():
                folder_cache[f.name.lower().strip()] = f

            created_folders: list[str] = []
            color_idx = 0

            for fi in files_to_categorize:
                if _wait_while_paused(job_id):
                    _update_job(job_id, status="cancelled")
                    return

                fname = fi["name"]
                suggested_folder = full_mapping.get(fname, "Uncategorized")

                # TRASH → "Review Later" folder (safe, not deleted)
                is_trash = suggested_folder.upper() == "TRASH"
                if is_trash:
                    suggested_folder = REVIEW_FOLDER_NAME

                # Find or create the folder
                folder_key = suggested_folder.lower().strip()
                if folder_key in folder_cache:
                    folder_obj = folder_cache[folder_key]
                else:
                    # Create new folder
                    color = AUTO_FOLDER_COLORS[color_idx % len(AUTO_FOLDER_COLORS)]
                    color_idx += 1
                    description = "AI-suggested non-educational files for review" if is_trash else "Auto-created from Smart Import"
                    folder_obj = Folder(
                        id=str(uuid.uuid4()),
                        name=suggested_folder.strip(),
                        description=description,
                        color="hsl(0 0% 45%)" if is_trash else color,  # grey for Review Later
                    )
                    db.session.add(folder_obj)
                    db.session.flush()  # get the id immediately
                    folder_cache[folder_key] = folder_obj
                    created_folders.append(suggested_folder.strip())

                # Copy file to uploads directory
                try:
                    unique_name = f"{uuid.uuid4().hex[:8]}_{fname}"
                    dest_path = os.path.join(upload_folder, unique_name)
                    shutil.copy2(fi["path"], dest_path)
                    fsize = os.path.getsize(dest_path)

                    record = UploadedFileRecord(
                        id=str(uuid.uuid4()),
                        folder_id=folder_obj.id,
                        original_name=fname,
                        file_size=fsize,
                        file_type=fi["ext"],
                        input_mode="files",
                        stored_path=dest_path,
                        processing_status="completed" if is_trash else "pending"
                    )
                    db.session.add(record)

                    if is_trash:
                        _update_file_status(job_id, fname, "review",
                                            folder_name=REVIEW_FOLDER_NAME,
                                            reason="AI suggests this may not be educational material")
                    else:
                        _update_file_status(job_id, fname, "done", folder_name=suggested_folder.strip())

                except Exception as e:
                    logger.warning("Failed to import file %s: %s", fname, e)
                    _update_file_status(job_id, fname, "error", reason=str(e)[:200])

            db.session.commit()

            # Trigger document processing for imported records (skip "Review Later" files)
            _trigger_background_processing(app, job_id)

            _update_job(job_id, status="completed", createdFolders=created_folders)
            logger.info("Smart import job %s completed: %d files processed", job_id, len(files_to_categorize))

        except Exception as e:
            logger.error("Smart import job %s failed: %s", job_id, e, exc_info=True)
            _update_job(job_id, status="error", error=str(e)[:500])
            try:
                db.session.rollback()
            except Exception:
                pass


def _trigger_background_processing(app, job_id: str):
    """Trigger document processing for all successfully imported records (not review).
    Runs in a SINGLE background thread to avoid ChromaDB concurrent init issues.
    """
    job = get_job(job_id)
    if not job:
        return

    done_files = [f for f in job["files"] if f["status"] == "done"]
    if not done_files:
        return

    def _process_all():
        with app.app_context():
            for fi in done_files:
                try:
                    record = UploadedFileRecord.query.filter_by(
                        original_name=fi["name"]
                    ).order_by(
                        UploadedFileRecord.created_at.desc()
                    ).first()
                    if record:
                        from app.features.upload.document_processor import process_record
                        process_record(record.id)
                except Exception as e:
                    logger.error("Background processing failed for file %s: %s", fi["name"], e)

    t = threading.Thread(target=_process_all, daemon=True)
    t.start()
